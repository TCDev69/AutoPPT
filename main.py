import openai
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx import *
from pptx.enum.text import PP_ALIGN

openai.api_key = "sk-KEY_HERE"
model_engine = "gpt-3.5-turbo" 

prs = Presentation()

#set slide template    
TITLE = 0
CONTENT = 1

#input
Titolo = input("Inserisci il titolo della presentazione: ")
Crediti = input("inserisci i crediti, es. fatto da X: ")
Sld = range(int(input("Inserisci il numero delle slide: ")))

title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = Titolo
subtitle.text = Crediti
title.alignment = PP_ALIGN.CENTER
subtitle.alignment = PP_ALIGN.CENTER
title.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
subtitle.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

for i in Sld:
    Arg1 = input("Inserisci il TITOLO del'argomento: ")
    ## CHATGPT
    response = openai.ChatCompletion.create(
        model='gpt-3.5-turbo',
        messages=[
            {"role": "system", "content": "dovrai scrivere un testo di tante righe, in un solo paragrafo"},
            {"role": "user", "content": "scrivi un testo in tante parole su: " + Arg1},
        ])

    message = response.choices[0]['message']
    print("Argomento 1\n")
    print("{}".format(message['content']))
    
    #Set title page text variables
    slide1_title = Arg1

    #Add a title slide
    slide1_layout = prs.slide_layouts[CONTENT]
    slide1 = prs.slides.add_slide(slide1_layout)
    slide1.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    slide1.shapes.title.text = slide1_title
    textbox1 = slide1.shapes.add_textbox(Inches(0.7), Inches(1.5),Inches(12), Inches(5.8))
    textbox1.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    textframe1 = textbox1.text_frame
    paragraph1 = textframe1.add_paragraph()
    paragraph1.text = "{}".format(message['content'])
    textframe1.margin_bottom = Inches(0)
    textframe1.margin_left = Inches(0.1)
    textframe1.vertical_anchor = MSO_ANCHOR.TOP
    textframe1.word_wrap = True
    textframe1.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    prs.slide_width = 11887200
    prs.slide_height = 6686550
    

#Salva

print("finito!")

prs.save("file.pptx")

print("file salvato!")
